from pptx import Presentation
from pptx.util import Emu

prs = Presentation("./templates/firm-template.pptx")

print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU "
      f"({Emu(prs.slide_width).inches:.2f}in x {Emu(prs.slide_height).inches:.2f}in)")

print("\n=== SLIDE LAYOUTS ===")
for i, layout in enumerate(prs.slide_layouts):
    print(f"[{i}] {layout.name}")
    for ph in layout.placeholders:
        print(f"    placeholder idx={ph.placeholder_format.idx} type={ph.placeholder_format.type} name={ph.name}")

print("\n=== EXISTING SLIDES ===")
for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    print(f"\n--- Slide {i+1}: layout='{layout_name}' ---")
    for shape in slide.shapes:
        kind = "placeholder" if shape.is_placeholder else shape.shape_type
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text[:80].replace("\n", " | ")
        print(f"  shape name={shape.name!r} kind={kind} pos=({Emu(shape.left).inches:.2f},{Emu(shape.top).inches:.2f}) "
              f"size=({Emu(shape.width).inches:.2f}x{Emu(shape.height).inches:.2f}) text={text!r}")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    color = None
                    try:
                        color = font.color.rgb
                    except Exception:
                        pass
                    print(f"      run text={run.text[:40]!r} font={font.name} size={font.size} bold={font.bold} color={color}")
