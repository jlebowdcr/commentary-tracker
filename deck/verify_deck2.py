import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Emu

prs = Presentation("./out/MELI_commentary_digest.pptx")
LOGO_LEFT_IN = 8.33

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    empty_boxes = 0
    for shape in slide.shapes:
        text = shape.text_frame.text if shape.has_text_frame else None
        if shape.is_placeholder and (text is None or not text.strip()):
            empty_boxes += 1
            print(f"  EMPTY PLACEHOLDER: idx={shape.placeholder_format.idx} name={shape.name}")
        if shape == slide.shapes.title if slide.shapes.title else False:
            right_edge = Emu(shape.left + shape.width).inches
            print(f"  TITLE right edge = {right_edge:.2f}in (logo starts at {LOGO_LEFT_IN}in) "
                  f"{'OVERLAP!' if right_edge > LOGO_LEFT_IN else 'ok'}")
    if empty_boxes == 0:
        print("  no empty placeholders")
