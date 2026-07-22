from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x00, 0x3C, 0x60)
BLACK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xFF, 0x00, 0x00)

LOGO_PATH = "./templates/logo.png"
LOGO_POS = dict(left=Inches(8.33), top=Inches(0.15), width=Inches(1.18), height=Inches(1.11))


def delete_all_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        prs.part.drop_rel(sld_id.rId)
        xml_slides.remove(sld_id)


def add_logo(slide, prs):
    slide.shapes.add_picture(LOGO_PATH, **LOGO_POS)


def strip_unused_placeholders(slide, keep_idx=(0,)):
    """Layouts in this template define extra empty body placeholders (idx 1 and 2
    on 'Title and Content') that we don't use — remove them so they don't render
    as empty click-to-add-text boxes."""
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx not in keep_idx:
            ph._element.getparent().remove(ph._element)


def fit_title(slide, width=7.6, height=1.1):
    """Narrow/heighten the title box so long headline-style titles wrap instead
    of running underneath the logo (logo starts at left=8.33in)."""
    title = slide.shapes.title
    title.width = Inches(width)
    title.height = Inches(height)


def add_bullets(slide, bullets, top=1.7, height=4.1, size=12, bold=True):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9.0), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = BLACK
        p.space_after = Pt(10)
        # blank spacer paragraph, matching template convention
        spacer = tf.add_paragraph()
        spacer.add_run().text = ""
    return box


def add_citation(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(9.0), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.bold = False
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_content_slide(prs, title, bullets, citation=None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    strip_unused_placeholders(slide)
    slide.shapes.title.text = title
    fit_title(slide)
    add_bullets(slide, bullets)
    if citation:
        add_citation(slide, citation)
    add_logo(slide, prs)
    return slide


prs = Presentation("./templates/firm-template.pptx")
delete_all_slides(prs)

# --- Title slide ---
title_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_layout)
title_ph = slide.placeholders[0]
title_ph.text = "MercadoLibre (MELI)"
run = title_ph.text_frame.paragraphs[0].runs[0]
run.font.size = Pt(36)
run.font.bold = True
run.font.color.rgb = NAVY

subtitle_ph = slide.placeholders[1]
tf = subtitle_ph.text_frame
tf.text = "Commentary Digest — YouTube & Podcast Sources"
tf.paragraphs[0].runs[0].font.size = Pt(18)
tf.paragraphs[0].runs[0].font.bold = True
for line in ["July 2026", "Delta Capital Research"]:
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.size = Pt(16)
    r.font.bold = True
slide.shapes.add_picture(LOGO_PATH, left=Inches(8.33), top=Inches(0.15), width=Inches(1.18), height=Inches(1.11))

# --- Scope & sources ---
add_content_slide(
    prs,
    "Scope: YouTube and podcast sources only — Reddit and X not yet available",
    [
        "Target: MercadoLibre (MELI); window: recent commentary (2025–2026)",
        "Sources queried: YouTube (3 targeted searches), podcast RSS (Business Breakdowns feed)",
        "Reddit: API access requested 2026-07-06, pending Reddit's manual review",
        "X: on hold — no free API tier since Feb 2026 (pay-per-use only)",
    ],
    citation="Full source-level detail: MELI_commentary_digest_2026-07-06.md",
)

# --- Context ---
add_content_slide(
    prs,
    "Leadership transition just completed: Galperin to Szarfsztejn",
    [
        "Marcos Galperin stepped down as CEO Jan 1, 2026 after ~25 years; became Executive Chairman",
        "Ariel Szarfsztejn (joined 2017; built logistics org from 60 to 70,000+ employees) is now CEO",
        "Galperin gave up his salary and title specifically to force full commitment onto the new CEO",
    ],
)

# --- Succession & leadership philosophy ---
add_content_slide(
    prs,
    "Galperin engineered the succession as a forcing function, not a retirement",
    [
        '"I’m giving up my salary, I’m giving up my job, I’m giving up the position '
        'because whoever leads this company cannot be focused on Mercado Libre and other things." '
        "— Galperin",
        'Szarfsztejn calls it Galperin’s "visionary move ... a generational change" — says '
        "he was already acting as CEO before the official date",
        "Watch item: Galperin says he remains involved ‘24/7, 365’ as Executive Chairman — "
        "worth tracking whether Szarfsztejn has real autonomy",
    ],
    citation="Bain & Company, “Founder’s Mentality: The CEO Sessions” (2026-02-24) · "
    "Americas Society/COA, “BRAVO Leadership Conversation” (2025-10-16)",
)

# --- Competitive positioning ---
add_content_slide(
    prs,
    "Amazon's 2015 entry gets credit for sharpening MELI; Chinese rivals seen as a tougher fight",
    [
        'Galperin: "we would not be the Mercado Libre that we are today if it were not because '
        'of Amazon coming into Mexico [in 2015]"',
        "Chinese competitors (Temu/Shein) described as less price-disciplined than Amazon — "
        '“a long and tough battle”',
        'Szarfsztejn: Brazil is "probably the most competitive scenario in e-commerce in the world"',
        "Independent analyst (Sven Carlin, June 2026) flags Wall Street concern over margin "
        "suppression tied to competitive pressure",
    ],
)

# --- AI strategy ---
add_content_slide(
    prs,
    "AI is already capping developer headcount growth, not just an efficiency story",
    [
        "Szarfsztejn's four pillars: efficiency, developer productivity, product improvements, "
        "early-stage agentic experience (piloted inside Mercado Pago)",
        'Galperin: "this is the first year ever that we won’t grow our developers" — projects '
        "20,000 to ~10,000 over five years via attrition, not layoffs",
        'Platform risk flagged by the founder directly: an AI-native device/partnership could mean '
        '"our app doesn’t show up anymore"',
        "Gap: no independent pushback found yet on this narrative — check once Reddit/X access "
        "is available",
    ],
)

# --- Fintech / credit inclusion ---
add_content_slide(
    prs,
    "Management frames fintech, not e-commerce, as the bigger growth lever",
    [
        "Mexico: ~6% loan-to-GDP ratio; ~85% of Mexicans lack a credit card — framed as the "
        "scale of the opportunity",
        "MELI credit is often cited as a small business's first-ever access to credit",
        "Argentina's remunerated-deposit product grew from ~400K to 18 million investors on-platform",
    ],
)

# --- Independent investor take ---
add_content_slide(
    prs,
    "Independent analyst sees fair value near current levels, not a clear mispricing",
    [
        "Sven Carlin (independent, YouTube, 2026-06-07): 42% revenue growth, 47% growth in items "
        "sold last quarter",
        "His own DCF-style estimate: ~$1,000/share conservative intrinsic value; implies 10–15% "
        "expected return depending on growth assumptions",
        'Claims Michael Burry holds a position and "is waiting for lower prices" — [UNSOURCED], '
        "the analyst's own claim, not independently verified",
    ],
)

# --- Coverage gaps ---
add_content_slide(
    prs,
    "Two sources still closed; one podcast transcript blocked",
    [
        "Reddit — application submitted 2026-07-06; approval odds uncertain for a private "
        "research use case under the new Responsible Builder Policy",
        "X — on hold; pay-per-use only since Feb 2026 (no free tier)",
        "Business Breakdowns podcast episode on MELI (Sept 2025): show notes retrieved; full "
        "transcript blocked by source site (403), no transcription key configured",
        "Two Spanish-language YouTube videos returned no transcript (caption-language gap)",
    ],
)

prs.save("./out/MELI_commentary_digest.pptx")
print("Saved ./out/MELI_commentary_digest.pptx")
print(f"Slide count: {len(prs.slides.__iter__.__self__._sldIdLst)}")
