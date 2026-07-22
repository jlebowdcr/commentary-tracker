from pptx import Presentation
from pptx.util import Emu

prs = Presentation("./templates/firm-template.pptx")

layout = prs.slide_layouts[1]  # "Title and Content"
print(f"Layout: {layout.name}")
print(f"Total shapes (incl. placeholders): {len(layout.shapes)}")
for shape in layout.shapes:
    is_ph = shape.is_placeholder
    print(f"  name={shape.name!r} is_placeholder={is_ph} kind={shape.shape_type} "
          f"pos=({Emu(shape.left).inches:.2f},{Emu(shape.top).inches:.2f}) "
          f"size=({Emu(shape.width).inches:.2f}x{Emu(shape.height).inches:.2f})")

# also check theme colors on the master
master = prs.slide_masters[0]
print("\nTheme color scheme:")
theme_el = master.element.getroottree().getroot()
