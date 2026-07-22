from pptx import Presentation
from pptx.util import Emu

prs = Presentation("./templates/firm-template.pptx")
slide2 = prs.slides[1]  # slide index 1 = "Slide 2" from inspection
for shape in slide2.shapes:
    if shape.shape_type == 13:  # PICTURE
        left_in = Emu(shape.left).inches
        top_in = Emu(shape.top).inches
        if abs(left_in - 8.33) < 0.05 and abs(top_in - 0.15) < 0.05:
            image = shape.image
            ext = image.ext
            with open(f"./templates/logo.{ext}", "wb") as f:
                f.write(image.blob)
            print(f"Saved logo as ./templates/logo.{ext}, size={image.size}")

# Also grab the title-slide hero picture for reference (not required but useful)
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.shape_type == 13:
        image = shape.image
        with open(f"./templates/title_hero.{image.ext}", "wb") as f:
            f.write(image.blob)
        print(f"Saved title hero as ./templates/title_hero.{image.ext}")
