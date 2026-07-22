from pptx import Presentation
from pptx.util import Emu

prs = Presentation("./out/MELI_commentary_digest.pptx")
print(f"Slide count: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} (layout={slide.slide_layout.name}) ---")
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f"  [{shape.shape_type}] {shape.text_frame.text[:100]!r}")
        elif shape.shape_type == 13:
            print(f"  [PICTURE] pos=({Emu(shape.left).inches:.2f},{Emu(shape.top).inches:.2f})")
